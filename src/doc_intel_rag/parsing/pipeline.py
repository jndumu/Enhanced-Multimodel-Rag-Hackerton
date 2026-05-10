"""GLM-OCR + PP-DocLayout-V3 document parsing pipeline."""

from __future__ import annotations

import asyncio
import base64
import hashlib
import tempfile
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any

from loguru import logger
from tenacity import retry, retry_if_exception_type, stop_after_attempt, wait_exponential

from doc_intel_rag.config import Settings, get_settings
from doc_intel_rag.parsing.entity_types import EntityLabel, ENTITY_TO_MODALITY


@dataclass
class BBox:
    """Axis-aligned bounding box of a detected layout element.

    Coordinates are in PDF points (72 dpi) relative to the page origin
    (bottom-left for most PDF renderers).

    Attributes:
        x0: Left edge of the bounding box.
        y0: Bottom edge of the bounding box.
        x1: Right edge of the bounding box.
        y1: Top edge of the bounding box.
        page: 1-based page number on which this box appears.
    """

    x0: float
    y0: float
    x1: float
    y1: float
    page: int


@dataclass
class ParsedElement:
    """A single layout element detected by PP-DocLayout-V3 / GLM-OCR.

    Attributes:
        label: Canonical :class:`~doc_intel_rag.parsing.entity_types.EntityLabel`.
        text: Extracted text content (OCR output or raw text for PDFs).
        page: 1-based page number.
        confidence: Layout detection confidence in ``[0, 1]``.
        bbox: Bounding box of the element on the page; ``None`` if unavailable.
        raw_image_b64: Base64-encoded PNG crop of the element's bounding region;
            populated for image, graph, and table elements.
        latex: LaTeX string for formula elements; ``None`` for all others.
        html: HTML representation for table elements; ``None`` for all others.
    """

    label: EntityLabel
    text: str
    page: int
    confidence: float
    bbox: BBox | None = None
    raw_image_b64: str | None = None
    latex: str | None = None
    html: str | None = None

    @property
    def modality(self) -> str:
        """Lookup the canonical modality string for this element's label."""
        return ENTITY_TO_MODALITY[self.label]


@dataclass
class ParseResult:
    """Complete output of a single document parse pass.

    Attributes:
        doc_id: SHA-256 hex digest of the source file bytes.  Used as the
            idempotency key for ingestion.
        source_file: Absolute path or URL of the document that was parsed.
        page_count: Total number of pages in the document.
        elements: Ordered list of all detected :class:`ParsedElement` objects.
        raw_metadata: Provider-specific metadata returned by the GLM-OCR SDK.
    """

    doc_id: str
    source_file: str
    page_count: int
    elements: list[ParsedElement] = field(default_factory=list)
    raw_metadata: dict[str, Any] = field(default_factory=dict)


class DocumentParser:
    """Entry point for document parsing using GLM-OCR + PP-DocLayout-V3.

    Supports both ``cloud`` and ``local`` backends as configured by
    ``settings.glmocr_backend``.  When the ``glmocr`` package is not installed
    or the API key is absent, a :class:`_StubGLMOCRClient` is used so the rest
    of the pipeline can still be exercised in tests without the SDK.

    All API calls are retried up to 3 times with exponential back-off via
    ``tenacity``.  CPU-bound operations (page rendering, bbox crop) are
    offloaded to a thread pool via ``asyncio.get_event_loop().run_in_executor``.

    Args:
        settings: Runtime configuration. Defaults to the global singleton.
    """

    def __init__(self, settings: Settings | None = None) -> None:
        self._settings = settings or get_settings()
        self._client: Any = None

    def _get_client(self) -> Any:
        if self._client is not None:
            return self._client

        try:
            import glmocr  # type: ignore[import-untyped]

            if self._settings.glmocr_backend == "cloud":
                self._client = glmocr.Client(
                    api_key=self._settings.glmocr_api_key,
                    timeout=self._settings.glmocr_timeout,
                )
            else:
                self._client = glmocr.LocalClient(
                    model=self._settings.glmocr_local_model,
                    timeout=self._settings.glmocr_timeout,
                )
        except ImportError:
            logger.info("glmocr SDK not installed — using PyMuPDF fallback parser")
            self._client = _PyMuPDFClient()

        return self._client

    @retry(
        retry=retry_if_exception_type(Exception),
        wait=wait_exponential(multiplier=1, min=2, max=30),
        stop=stop_after_attempt(3),
        reraise=True,
    )
    async def parse(self, file_path: str | Path) -> ParseResult:
        raw_bytes, resolved_path = await self._load_source(file_path)
        doc_id = hashlib.sha256(raw_bytes).hexdigest()

        logger.info("Parsing document", doc_id=doc_id[:12], path=resolved_path)

        client = self._get_client()

        # When source is a URL the client receives a local temp file so that
        # extension-based dispatch (PDF vs DOCX etc.) works correctly.
        # Detect format from magic bytes so numeric URL paths (e.g. arXiv /pdf/2304.12306)
        # are still treated as PDF.
        loop = asyncio.get_running_loop()
        if resolved_path.startswith(("http://", "https://")):
            if raw_bytes[:4] == b"%PDF":
                suffix = ".pdf"
            elif raw_bytes[:2] == b"PK":
                url_suffix = Path(resolved_path.split("?")[0]).suffix.lower()
                suffix = url_suffix if url_suffix in {".pptx", ".xlsx"} else ".docx"
            else:
                suffix = Path(resolved_path.split("?")[0]).suffix or ".txt"
            with tempfile.NamedTemporaryFile(suffix=suffix, delete=False) as tmp:
                tmp.write(raw_bytes)
                local_path = tmp.name
            raw_result = await loop.run_in_executor(
                None, lambda: client.parse(local_path)
            )
            Path(local_path).unlink(missing_ok=True)
        else:
            raw_result = await loop.run_in_executor(
                None, lambda: client.parse(resolved_path)
            )

        elements = self._convert_elements(raw_result, raw_bytes)

        # OCR pass: find scanned/image-heavy pages and extract text via vision model
        if resolved_path.endswith(".pdf") and self._settings.vision_enabled:
            elements = await self._ocr_scanned_pages(raw_bytes, elements)

        page_count = max((e.page for e in elements), default=1)

        logger.info(
            "Parsing complete",
            doc_id=doc_id[:12],
            elements=len(elements),
            pages=page_count,
        )

        return ParseResult(
            doc_id=doc_id,
            source_file=resolved_path,
            page_count=page_count,
            elements=elements,
            raw_metadata=getattr(raw_result, "metadata", {}),
        )

    async def _ocr_scanned_pages(
        self, raw_bytes: bytes, elements: list[ParsedElement]
    ) -> list[ParsedElement]:
        """Run vision-model OCR on PDF pages that have images but no extracted text.

        Called automatically for PDF files when ``vision_enabled=True``.
        Pages that already have text content (from PyMuPDF) are skipped.
        Pages that are mostly images (scanned documents) are rendered to PNG
        and sent to the configured ``vision_model`` (default: ``llava``) via
        the Ollama OpenAI-compatible endpoint.
        """
        try:
            import fitz
        except ImportError:
            return elements

        pages_with_text: set[int] = {
            e.page for e in elements if e.text.strip() and len(e.text) > 30
        }

        try:
            doc = fitz.open(stream=raw_bytes, filetype="pdf")
        except Exception:
            return elements

        new_elements: list[ParsedElement] = list(elements)

        for page_num_zero, page in enumerate(doc):
            page_num = page_num_zero + 1
            if page_num in pages_with_text:
                continue  # already has text

            # Only process pages that actually have images
            if not page.get_images(full=False):
                continue

            try:
                pix = page.get_pixmap(dpi=200)
                img_b64 = base64.b64encode(pix.tobytes("png")).decode()

                ocr_text = await self._vision_ocr(img_b64, page_num)
                if ocr_text and len(ocr_text.strip()) > 10:
                    new_elements.append(ParsedElement(
                        label=EntityLabel.PARAGRAPH,
                        text=ocr_text.strip(),
                        page=page_num,
                        confidence=0.75,
                        raw_image_b64=img_b64,
                    ))
                    logger.info("OCR extracted text from page", page=page_num, chars=len(ocr_text))
            except Exception as exc:
                logger.debug("OCR failed for page", page=page_num, error=str(exc))

        doc.close()
        return new_elements

    async def _vision_ocr(self, img_b64: str, page: int) -> str:
        """Send a page image to the vision model and return extracted text."""
        try:
            from openai import AsyncOpenAI

            client = AsyncOpenAI(
                api_key=self._settings.mesh_api_key,
                base_url=self._settings.mesh_api_base_url,
            )
            response = await client.chat.completions.create(
                model=self._settings.vision_model,
                messages=[{
                    "role": "user",
                    "content": [
                        {
                            "type": "text",
                            "text": (
                                "Read this document page and transcribe every word of text "
                                "you can see, in reading order (top to bottom, left to right). "
                                "Include headings, body text, captions, and labels. "
                                "Do not describe the image — output the raw text only."
                            ),
                        },
                        {
                            "type": "image_url",
                            "image_url": {"url": f"data:image/png;base64,{img_b64}"},
                        },
                    ],
                }],
                max_tokens=2048,
                temperature=0,
            )
            return response.choices[0].message.content or ""
        except Exception as exc:
            logger.debug("Vision OCR call failed", page=page, error=str(exc))
            return ""

    async def _load_source(self, file_path: str | Path) -> tuple[bytes, str]:
        """Load raw bytes from a local path or HTTP/HTTPS URL.

        Returns:
            Tuple of (raw_bytes, resolved_path_string).
        """
        source = str(file_path)
        if source.startswith(("http://", "https://")):
            import httpx
            async with httpx.AsyncClient(timeout=60, follow_redirects=True) as client:
                response = await client.get(source)
                response.raise_for_status()
                return response.content, source
        path = Path(file_path)
        if not path.exists():
            raise FileNotFoundError(f"Document not found: {path}")
        return path.read_bytes(), str(path)

    def _convert_elements(self, raw_result: Any, raw_bytes: bytes) -> list[ParsedElement]:
        elements: list[ParsedElement] = []

        raw_elements: list[Any] = getattr(raw_result, "elements", raw_result) or []

        for raw in raw_elements:
            label_str: str = getattr(raw, "label", "paragraph")
            try:
                label = EntityLabel(label_str)
            except ValueError:
                logger.debug("Unknown entity label, defaulting to paragraph", label=label_str)
                label = EntityLabel.PARAGRAPH

            bbox_raw = getattr(raw, "bbox", None)
            bbox: BBox | None = None
            if bbox_raw is not None:
                bbox = BBox(
                    x0=float(getattr(bbox_raw, "x0", 0)),
                    y0=float(getattr(bbox_raw, "y0", 0)),
                    x1=float(getattr(bbox_raw, "x1", 0)),
                    y1=float(getattr(bbox_raw, "y1", 0)),
                    page=int(getattr(raw, "page", 1)),
                )

            image_b64: str | None = None
            if ENTITY_TO_MODALITY[label] in {"image", "graph", "table"} and bbox is not None:
                image_b64 = self._crop_element(raw_bytes, bbox, str(getattr(raw, "page", 1)))

            elements.append(
                ParsedElement(
                    label=label,
                    text=str(getattr(raw, "text", "") or ""),
                    page=int(getattr(raw, "page", 1)),
                    confidence=float(getattr(raw, "confidence", 1.0)),
                    bbox=bbox,
                    raw_image_b64=image_b64 or getattr(raw, "image_b64", None),
                    latex=getattr(raw, "latex", None),
                    html=getattr(raw, "html", None),
                )
            )

        return elements

    def _crop_element(self, raw_bytes: bytes, bbox: BBox, page_str: str) -> str | None:
        try:
            import fitz  # PyMuPDF

            doc = fitz.open(stream=raw_bytes, filetype="pdf")
            page_num = max(0, int(page_str) - 1)
            if page_num >= len(doc):
                return None
            page = doc[page_num]
            rect = fitz.Rect(bbox.x0, bbox.y0, bbox.x1, bbox.y1)
            pix = page.get_pixmap(clip=rect, dpi=150)
            return base64.b64encode(pix.tobytes("png")).decode()
        except Exception as exc:
            logger.debug("Element crop failed", error=str(exc))
            return None


class _PyMuPDFClient:
    """Real document parser using PyMuPDF + python-docx + markdown.

    Handles: PDF, DOCX, PPTX, HTML, Markdown, plain text.
    Falls back gracefully to raw text extraction when format is unknown.

    This is the production fallback when the GLM-OCR SDK is not installed.
    It produces real :class:`_ParsedElement` objects with correct page numbers,
    basic entity label classification (title vs paragraph), and bounding boxes.
    """

    def parse(self, path: str) -> "_FallbackResult":
        suffix = Path(path).suffix.lower()
        try:
            if suffix == ".pdf":
                return self._parse_pdf(path)
            elif suffix in {".docx", ".doc"}:
                return self._parse_docx(path)
            elif suffix in {".pptx", ".ppt"}:
                return self._parse_pptx(path)
            elif suffix in {".md", ".markdown"}:
                return self._parse_markdown(path)
            elif suffix in {".html", ".htm"}:
                return self._parse_html(path)
            else:
                return self._parse_text(path)
        except Exception as exc:
            logger.warning("Fallback parser failed", path=path, error=str(exc))
            return _FallbackResult([_FallbackElement(
                label="paragraph",
                text=f"[Parse error: {exc}]",
                page=1, confidence=0.0,
            )])

    def _parse_pdf(self, path: str) -> "_FallbackResult":
        import fitz  # PyMuPDF

        elements: list[_FallbackElement] = []
        doc = fitz.open(path)

        for page_num, page in enumerate(doc, start=1):
            blocks = page.get_text("dict", flags=fitz.TEXT_PRESERVE_WHITESPACE)["blocks"]

            for block in blocks:
                if block.get("type") != 0:  # 0 = text, 1 = image
                    continue

                for line in block.get("lines", []):
                    for span in line.get("spans", []):
                        text = span.get("text", "").strip()
                        if not text:
                            continue

                        font_size: float = span.get("size", 12)
                        flags: int = span.get("flags", 0)
                        is_bold = bool(flags & 2**4)

                        # Classify by font size + bold
                        if font_size >= 16 or (font_size >= 13 and is_bold):
                            label = "document_title" if page_num == 1 else "section_title"
                        elif font_size >= 12 and is_bold:
                            label = "subsection_title"
                        else:
                            label = "paragraph"

                        r = span.get("bbox", (0, 0, 0, 0))
                        elements.append(_FallbackElement(
                            label=label, text=text, page=page_num, confidence=1.0,
                            bbox=BBox(x0=r[0], y0=r[1], x1=r[2], y1=r[3], page=page_num),
                        ))

        doc.close()
        logger.info("PDF parsed via PyMuPDF", path=path, elements=len(elements))
        return _FallbackResult(elements)

    def _parse_docx(self, path: str) -> "_FallbackResult":
        try:
            import docx  # python-docx
        except ImportError:
            return self._parse_text(path)

        elements: list[_FallbackElement] = []
        doc = docx.Document(path)
        page = 1

        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                continue

            style_name = (para.style.name or "").lower()
            if "heading 1" in style_name or "title" in style_name:
                label = "section_title"
            elif "heading" in style_name:
                label = "subsection_title"
            else:
                label = "paragraph"

            elements.append(_FallbackElement(label=label, text=text, page=page, confidence=1.0))

        for table in doc.tables:
            rows = [[cell.text.strip() for cell in row.cells] for row in table.rows]
            html = "<table>" + "".join(
                "<tr>" + "".join(f"<td>{c}</td>" for c in row) + "</tr>"
                for row in rows
            ) + "</table>"
            flat = " | ".join(cell for row in rows for cell in row if cell)
            if flat:
                elements.append(_FallbackElement(
                    label="table", text=flat, page=page, confidence=1.0, html=html,
                ))

        logger.info("DOCX parsed", path=path, elements=len(elements))
        return _FallbackResult(elements)

    def _parse_pptx(self, path: str) -> "_FallbackResult":
        try:
            from pptx import Presentation  # python-pptx
        except ImportError:
            return self._parse_text(path)

        elements: list[_FallbackElement] = []
        prs = Presentation(path)

        for slide_num, slide in enumerate(prs.slides, start=1):
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if not text:
                        continue
                    level = para.level
                    label = "section_title" if level == 0 else "paragraph"
                    elements.append(_FallbackElement(
                        label=label, text=text, page=slide_num, confidence=1.0,
                    ))

        logger.info("PPTX parsed", path=path, elements=len(elements))
        return _FallbackResult(elements)

    def _parse_markdown(self, path: str) -> "_FallbackResult":
        content = Path(path).read_text(encoding="utf-8", errors="replace")
        elements: list[_FallbackElement] = []
        page = 1

        for line in content.splitlines():
            stripped = line.strip()
            if not stripped:
                continue
            if stripped.startswith("# "):
                elements.append(_FallbackElement("document_title", stripped[2:], page, 1.0))
            elif stripped.startswith("## "):
                elements.append(_FallbackElement("section_title", stripped[3:], page, 1.0))
            elif stripped.startswith("### "):
                elements.append(_FallbackElement("subsection_title", stripped[4:], page, 1.0))
            elif stripped.startswith(("- ", "* ", "+ ")):
                elements.append(_FallbackElement("list_item", stripped[2:], page, 1.0))
            elif stripped.startswith("> "):
                elements.append(_FallbackElement("blockquote", stripped[2:], page, 1.0))
            else:
                elements.append(_FallbackElement("paragraph", stripped, page, 1.0))

        logger.info("Markdown parsed", path=path, elements=len(elements))
        return _FallbackResult(elements)

    def _parse_html(self, path: str) -> "_FallbackResult":
        content = Path(path).read_text(encoding="utf-8", errors="replace")
        try:
            from html.parser import HTMLParser

            class _Collector(HTMLParser):
                def __init__(self) -> None:
                    super().__init__()
                    self.items: list[tuple[str, str]] = []
                    self._tag = "paragraph"
                    self._buf: list[str] = []

                def handle_starttag(self, tag: str, attrs: object) -> None:
                    if tag in ("h1", "h2", "h3", "h4"):
                        self._flush()
                        self._tag = {"h1": "document_title", "h2": "section_title",
                                     "h3": "subsection_title", "h4": "subsection_title"}[tag]
                    elif tag in ("p", "li", "blockquote", "pre"):
                        self._flush()
                        self._tag = {"p": "paragraph", "li": "list_item",
                                     "blockquote": "blockquote", "pre": "code_block"}[tag]

                def handle_endtag(self, tag: str) -> None:
                    if tag in ("h1", "h2", "h3", "h4", "p", "li", "blockquote", "pre"):
                        self._flush()

                def handle_data(self, data: str) -> None:
                    self._buf.append(data)

                def _flush(self) -> None:
                    text = "".join(self._buf).strip()
                    if text:
                        self.items.append((self._tag, text))
                    self._buf = []
                    self._tag = "paragraph"

            collector = _Collector()
            collector.feed(content)
            collector._flush()
            elements = [
                _FallbackElement(label, text, 1, 1.0)
                for label, text in collector.items
            ]
        except Exception:
            elements = [_FallbackElement("paragraph", content[:4000], 1, 1.0)]

        logger.info("HTML parsed", path=path, elements=len(elements))
        return _FallbackResult(elements)

    def _parse_text(self, path: str) -> "_FallbackResult":
        try:
            content = Path(path).read_text(encoding="utf-8", errors="replace")
        except Exception:
            content = "[Could not read file]"
        elements = [
            _FallbackElement("paragraph", para.strip(), i + 1, 1.0)
            for i, para in enumerate(content.split("\n\n"))
            if para.strip()
        ]
        return _FallbackResult(elements)


@dataclass
class _FallbackResult:
    """Parse result from the PyMuPDF fallback parser."""
    _elements: list["_FallbackElement"]

    @property
    def elements(self) -> list[Any]:
        return self._elements

    @property
    def metadata(self) -> dict[str, Any]:
        return {"parser": "pymupdf_fallback"}


@dataclass
class _FallbackElement:
    label: str
    text: str
    page: int
    confidence: float
    bbox: "BBox | None" = None
    image_b64: None = None
    latex: None = None
    html: str | None = None
